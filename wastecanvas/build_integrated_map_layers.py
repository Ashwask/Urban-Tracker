from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


BASE_IMAGE = "/Users/ashwinkulkarni/Downloads/codex-map-preview/Attachment287.pdf.png"
OUT1 = "/Users/ashwinkulkarni/Documents/New project/india_map_layer_1_quantified.png"
OUT2 = "/Users/ashwinkulkarni/Documents/New project/india_map_layer_2_actors.png"
OUT3 = "/Users/ashwinkulkarni/Documents/New project/india_map_layer_3_capital.png"
PPT_OUT = "/Users/ashwinkulkarni/Documents/New project/Urban Waste System - India Integrated Map Layers.pptx"

BG = (247, 248, 244, 255)
TEXT = (34, 34, 34)
MUTED = (95, 95, 95)
BORDER = (206, 210, 204)
LIGHT = (255, 255, 255, 248)
RED = (206, 66, 66)
ORANGE = (230, 126, 34)
PURPLE = (155, 89, 182)
BLUE = (52, 152, 219)
GREEN = (46, 204, 113)
GRAY = (127, 140, 141)


def load_font(size: int, bold: bool = False):
    names = [
        "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    ]
    for name in names:
        try:
            return ImageFont.truetype(name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    line = ""
    for word in words:
        trial = word if not line else f"{line} {word}"
        if draw.textlength(trial, font=font) <= max_width:
            line = trial
        else:
            if line:
                lines.append(line)
            line = word
    if line:
        lines.append(line)
    return lines


def draw_box(
    draw: ImageDraw.ImageDraw,
    rect: tuple[int, int, int, int],
    radius: int,
    fill: tuple[int, int, int, int],
    outline: tuple[int, int, int, int],
    shadow: int = 6,
) -> None:
    x1, y1, x2, y2 = rect
    draw.rounded_rectangle([x1 + shadow, y1 + shadow, x2 + shadow, y2 + shadow], radius=radius, fill=(0, 0, 0, 18))
    draw.rounded_rectangle(rect, radius=radius, fill=fill, outline=outline, width=3)


def draw_arrow(draw: ImageDraw.ImageDraw, points: list[tuple[int, int]], color: tuple[int, int, int], width: int = 3) -> None:
    draw.line(points, fill=color + (255,), width=width, joint="curve")
    sx, sy = points[-2]
    ex, ey = points[-1]
    dx = ex - sx
    dy = ey - sy
    length = max((dx * dx + dy * dy) ** 0.5, 1)
    ux = dx / length
    uy = dy / length
    px = -uy
    py = ux
    head = 10
    wing = 6
    draw.polygon(
        [
            (ex, ey),
            (int(ex - head * ux + wing * px), int(ey - head * uy + wing * py)),
            (int(ex - head * ux - wing * px), int(ey - head * uy - wing * py)),
        ],
        fill=color + (255,),
    )


def connector_points(
    box: tuple[int, int, int, int],
    marker: tuple[int, int],
    side: str,
    map_rect: tuple[int, int, int, int],
) -> list[tuple[int, int]]:
    x, y, w, h = box
    mx, my = marker
    map_x1, map_y1, map_x2, map_y2 = map_rect

    if side == "right":
        start = (x + w, y + h // 2)
        elbow_x = map_x1 - 18
        return [start, (elbow_x, start[1]), (elbow_x, my), (mx - 18, my), marker]
    if side == "left":
        start = (x, y + h // 2)
        elbow_x = map_x2 + 18
        return [start, (elbow_x, start[1]), (elbow_x, my), (mx + 18, my), marker]
    if side == "bottom":
        start = (x + w // 2, y + h)
        elbow_y = map_y1 - 18
        return [start, (start[0], elbow_y), (mx, elbow_y), (mx, my - 18), marker]
    start = (x + w // 2, y)
    elbow_y = map_y2 + 18
    return [start, (start[0], elbow_y), (mx, elbow_y), (mx, my + 18), marker]


def scale_point(point: tuple[int, int], scale: float, offset: tuple[int, int]) -> tuple[int, int]:
    return (int(offset[0] + point[0] * scale), int(offset[1] + point[1] * scale))


def draw_callout(
    draw: ImageDraw.ImageDraw,
    fonts: dict[str, object],
    title: str,
    line1: str,
    line2: str,
    box: tuple[int, int, int, int],
    marker: tuple[int, int],
    side: str,
    color: tuple[int, int, int],
    map_rect: tuple[int, int, int, int],
) -> None:
    x, y, w, h = box
    draw_box(draw, (x, y, x + w, y + h), radius=18, fill=LIGHT, outline=color + (255,), shadow=4)
    draw.rounded_rectangle([x + 16, y + 14, x + 140, y + 24], radius=4, fill=color + (255,))
    draw.text((x + 16, y + 32), title, font=fonts["title"], fill=color + (255,))

    current_y = y + 60
    max_width = w - 28
    for line in wrap_text(draw, line1, fonts["line1"], max_width):
        draw.text((x + 14, current_y), line, font=fonts["line1"], fill=TEXT + (255,))
        current_y += 20
    current_y += 2
    for line in wrap_text(draw, line2, fonts["line2"], max_width):
        draw.text((x + 14, current_y), line, font=fonts["line2"], fill=MUTED + (255,))
        current_y += 18

    draw.ellipse([marker[0] - 9, marker[1] - 9, marker[0] + 9, marker[1] + 9], fill=color + (255,), outline=(255, 255, 255, 255), width=2)
    draw_arrow(draw, connector_points(box, marker, side, map_rect), color)


def make_canvas(title: str, subtitle: str):
    base = Image.open(BASE_IMAGE).convert("RGBA")
    canvas_w, canvas_h = 1920, 1080
    map_w = 1080
    scale = map_w / base.width
    map_h = int(base.height * scale)
    map_x, map_y = 420, 168
    map_rect = (map_x, map_y, map_x + map_w, map_y + map_h)

    canvas = Image.new("RGBA", (canvas_w, canvas_h), BG)
    draw = ImageDraw.Draw(canvas)
    fonts = {
        "header": load_font(30, bold=True),
        "sub": load_font(13),
        "title": load_font(15, bold=True),
        "line1": load_font(12, bold=True),
        "line2": load_font(11),
        "footer": load_font(10),
    }

    draw_box(draw, (22, 18, 1898, 96), radius=14, fill=(255, 255, 255, 246), outline=BORDER + (255,), shadow=0)
    draw.text((34, 28), title, font=fonts["header"], fill=TEXT + (255,))
    draw.text((36, 70), subtitle, font=fonts["sub"], fill=MUTED + (255,))

    draw_box(
        draw,
        (map_x - 8, map_y - 8, map_x + map_w + 8, map_y + map_h + 8),
        radius=16,
        fill=(255, 255, 255, 255),
        outline=BORDER + (255,),
        shadow=4,
    )
    resized = base.resize((map_w, map_h), Image.Resampling.LANCZOS)
    canvas.alpha_composite(resized, (map_x, map_y))
    return canvas, map_rect, scale, fonts


def save_footer(draw: ImageDraw.ImageDraw, fonts: dict[str, object], text: str) -> None:
    draw_box(draw, (22, 1018, 1898, 1060), radius=10, fill=(255, 255, 255, 240), outline=BORDER + (255,), shadow=0)
    draw.text((36, 1030), text, font=fonts["footer"], fill=MUTED + (255,))


def build_quantified() -> None:
    canvas, map_rect, scale, fonts = make_canvas(
        "Layer 1 | India-wide quantified ecosystem lens applied to the system map",
        "CPCB FY 2021-22 mass balance | GFC 2024-25 | PIB Jan/Mar 2026",
    )
    draw = ImageDraw.Draw(canvas)
    callouts = [
        {
            "title": "Generated + collected",
            "line1": "170,339 TPD generated | 156,449 TPD collected (92%)",
            "line2": "Collection scaled faster than material quality and routing.",
            "box": (26, 156, 348, 112),
            "anchor": (420, 805),
            "side": "right",
            "color": RED,
        },
        {
            "title": "Uncollected leakage",
            "line1": "13,890 TPD not collected",
            "line2": "Material still escapes before recovery even begins.",
            "box": (26, 286, 348, 106),
            "anchor": (700, 428),
            "side": "right",
            "color": RED,
        },
        {
            "title": "Source segregation gap",
            "line1": "60,478 wards >90% collection vs 39,648 >90% segregation",
            "line2": "This remains the biggest break in the chain.",
            "box": (26, 410, 348, 122),
            "anchor": (930, 620),
            "side": "right",
            "color": ORANGE,
        },
        {
            "title": "Post-collection leakage",
            "line1": "23,483 TPD collected but neither processed nor landfilled",
            "line2": "The system still loses material after collection.",
            "box": (26, 550, 348, 118),
            "anchor": (1190, 820),
            "side": "right",
            "color": PURPLE,
        },
        {
            "title": "Offtake bottleneck",
            "line1": "~35% of MSW is dry waste",
            "line2": "Sorting only matters if buyers, EPR and recyclers absorb it.",
            "box": (1546, 156, 348, 106),
            "anchor": (1700, 286),
            "side": "left",
            "color": PURPLE,
        },
        {
            "title": "MRF bottleneck",
            "line1": "2,900 MRFs | 67,000 TPD capacity",
            "line2": "Feedstock quality, not only capacity, is the binding issue.",
            "box": (1546, 280, 348, 112),
            "anchor": (1330, 592),
            "side": "left",
            "color": BLUE,
        },
        {
            "title": "Processing gap",
            "line1": "91,511 TPD processed (54%)",
            "line2": "Recovery capacity still trails total waste generation.",
            "box": (1546, 410, 348, 106),
            "anchor": (1768, 720),
            "side": "left",
            "color": GREEN,
        },
        {
            "title": "Landfill + legacy drag",
            "line1": "41,455 TPD landfilled (24%) | ~25 crore MT legacy waste",
            "line2": "~9.5 crore MT still remains nationally.",
            "box": (1546, 534, 348, 116),
            "anchor": (1500, 915),
            "side": "left",
            "color": GRAY,
        },
    ]
    for item in callouts:
        marker = scale_point(item["anchor"], scale, (map_rect[0], map_rect[1]))
        draw_callout(draw, fonts, item["title"], item["line1"], item["line2"], item["box"], marker, item["side"], item["color"], map_rect)

    save_footer(
        draw,
        fonts,
        "This layer reads the map as a national system problem: segregation quality, traceability, offtake, processing depth and legacy waste are more binding than first-mile collection alone.",
    )
    canvas.save(OUT1)


def build_actors() -> None:
    canvas, map_rect, scale, fonts = make_canvas(
        "Layer 2 | Specific actors and partners mapped onto each block of the India waste system",
        "Named examples across India; shown to widen the ecosystem lens, not to exclude local enterprises, nonprofits, cooperatives, consultants or family-backed actors.",
    )
    draw = ImageDraw.Draw(canvas)
    callouts = [
        {
            "title": "Consumers + bulk generators",
            "line1": "Indian Railways, DIAL, BIAL, Infosys, Taj Hotels, Apollo Hospitals",
            "line2": "These are the blocks where waste quality and source segregation are won or lost.",
            "box": (26, 150, 348, 112),
            "anchor": (340, 790),
            "side": "right",
            "color": PURPLE,
        },
        {
            "title": "Formal collectors",
            "line1": "Indore MC, Pune MC, BBMP, Antony Waste, Ramky Enviro, JBM",
            "line2": "Formal collection, routing and concessionaire execution sit here.",
            "box": (26, 278, 348, 112),
            "anchor": (760, 830),
            "side": "right",
            "color": BLUE,
        },
        {
            "title": "Informal collectors",
            "line1": "SWaCH, KKPKP, Hasiru Dala, Chintan, Safai Sena, Plastics For Change",
            "line2": "This block is essential to capture, livelihoods and citywide recovery.",
            "box": (26, 406, 348, 118),
            "anchor": (760, 315),
            "side": "right",
            "color": ORANGE,
        },
        {
            "title": "Waste transporters",
            "line1": "Antony Waste, Ramky Enviro, JBM, BBMP fleets, Pune SWaCH vehicles",
            "line2": "Transport links determine whether material reaches the right downstream node.",
            "box": (26, 540, 348, 118),
            "anchor": (1290, 910),
            "side": "right",
            "color": GREEN,
        },
        {
            "title": "Scrap dealers",
            "line1": "TheKabadiwala, Kabadiwalla Connect, Scrap Uncle, Paperman",
            "line2": "These traders remain the connective tissue of the dry-waste economy.",
            "box": (1546, 150, 348, 112),
            "anchor": (1195, 320),
            "side": "left",
            "color": PURPLE,
        },
        {
            "title": "Aggregators",
            "line1": "Recykal, GEM Enviro, Plastics For Change, Cirqlo, Nepra",
            "line2": "Aggregation platforms convert fragmented waste into bankable supply.",
            "box": (1546, 278, 348, 112),
            "anchor": (1670, 320),
            "side": "left",
            "color": BLUE,
        },
        {
            "title": "Recyclers",
            "line1": "Banyan Nation, Lucro, Attero, Gravita, Ecoreco, Deluxe Recycling",
            "line2": "These are the plants that turn collected material into circular output.",
            "box": (1546, 406, 348, 118),
            "anchor": (1890, 230),
            "side": "left",
            "color": ORANGE,
        },
        {
            "title": "MRFs / sorting centres",
            "line1": "Saahas Zero Waste, Hasiru Dala Innovations, Nepra, Chintan MRFs",
            "line2": "Sorting quality determines what becomes saleable and what becomes reject.",
            "box": (1546, 540, 348, 116),
            "anchor": (1365, 600),
            "side": "left",
            "color": GREEN,
        },
        {
            "title": "Waste processors",
            "line1": "GPS Renewables, IGRPL, Mailhem, Ramky Enviro, JITF Urban, Antony Lara",
            "line2": "Processing covers compost, biomethanation, CBG, RDF and scientific disposal.",
            "box": (420, 780, 1080, 130),
            "anchor": (1780, 705),
            "side": "top",
            "color": RED,
        },
    ]
    for item in callouts:
        marker = scale_point(item["anchor"], scale, (map_rect[0], map_rect[1]))
        draw_callout(draw, fonts, item["title"], item["line1"], item["line2"], item["box"], marker, item["side"], item["color"], map_rect)

    save_footer(
        draw,
        fonts,
        "These names are specific examples mapped to the blocks; the slide is intentionally additive and does not negate entrepreneurs, local operators, nonprofits, cooperatives, consultants or family-backed actors.",
    )
    canvas.save(OUT2)


def build_capital() -> None:
    canvas, map_rect, scale, fonts = make_canvas(
        "Layer 3 | Specific funders, philanthropy, CSR, investors and strategics attached to the same map",
        "Capital does not sit in one place. It enters through behavior change, inclusion, aggregation, recycling, processing, city systems and legacy remediation.",
    )
    draw = ImageDraw.Draw(canvas)
    callouts = [
        {
            "title": "Source segregation + behavior",
            "line1": "HUL-MoHUA, ITC WOW, Coca-Cola India Foundation Anandana, Tata Trusts",
            "line2": "BMGF, USAID and UNDP India also show up in behavior and systems work.",
            "box": (26, 150, 348, 122),
            "anchor": (420, 790),
            "side": "right",
            "color": PURPLE,
        },
        {
            "title": "Inclusion + livelihoods",
            "line1": "H&M Foundation, The Body Shop, Plastics For Change, Tata Trusts",
            "line2": "These names show up where informal workers are integrated into formal flows.",
            "box": (26, 290, 348, 116),
            "anchor": (760, 315),
            "side": "right",
            "color": ORANGE,
        },
        {
            "title": "City systems + public finance",
            "line1": "MoHUA, SBM-U 2.0, Smart Cities Mission, CPCB, BBMP, Pune MC",
            "line2": "Public money and public procurement shape first-mile, MRF and remediation choices.",
            "box": (26, 424, 348, 122),
            "anchor": (760, 830),
            "side": "right",
            "color": BLUE,
        },
        {
            "title": "Legacy remediation",
            "line1": "MoHUA, SBM-U 2.0, Lakshya Zero Dumpsite, Indore MC, Ahmedabad MC",
            "line2": "This is where public cleanup capital and execution are concentrated.",
            "box": (26, 564, 348, 112),
            "anchor": (1500, 915),
            "side": "right",
            "color": GRAY,
        },
        {
            "title": "MRF + traceability layer",
            "line1": "Social Alpha, Startup Gateway, Claypond Capital, Dr Ranjan Pai",
            "line2": "Vellayan Subbiah and Ajay Parekh show up through Recykal-linked growth capital.",
            "box": (1546, 150, 348, 122),
            "anchor": (1365, 600),
            "side": "left",
            "color": GREEN,
        },
        {
            "title": "Aggregation + recycling capital",
            "line1": "Circulate Capital, Morgan Stanley, Murugappa Group, Dow, Unilever",
            "line2": "The Coca-Cola Company and Recykal-linked investors attach here as well.",
            "box": (1546, 290, 348, 122),
            "anchor": (1670, 320),
            "side": "left",
            "color": BLUE,
        },
        {
            "title": "Corporate circularity backers",
            "line1": "PepsiCo, P&G, Danone, Chanel, Chevron Phillips, Mondelez",
            "line2": "These names show how corporate offtake and fund capital shape recycling markets.",
            "box": (1546, 430, 348, 118),
            "anchor": (1890, 230),
            "side": "left",
            "color": PURPLE,
        },
        {
            "title": "Processing + bio-CBG strategics",
            "line1": "Indian Oil, BPCL, Oil India, GPS Renewables, IGRPL, SATAT",
            "line2": "Strategic and infrastructure capital clusters around processing and fuel conversion.",
            "box": (1546, 566, 348, 116),
            "anchor": (1780, 705),
            "side": "left",
            "color": RED,
        },
    ]
    for item in callouts:
        marker = scale_point(item["anchor"], scale, (map_rect[0], map_rect[1]))
        draw_callout(draw, fonts, item["title"], item["line1"], item["line2"], item["box"], marker, item["side"], item["color"], map_rect)

    save_footer(
        draw,
        fonts,
        "Specific names are mapped to the flow as examples of philanthropy, CSR, investors, strategics and public finance. This is not an exclusion list.",
    )
    canvas.save(OUT3)


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for path in (OUT1, OUT2, OUT3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPT_OUT)


def main() -> None:
    build_quantified()
    build_actors()
    build_capital()
    build_ppt()


if __name__ == "__main__":
    main()
