from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT = "/Users/ashwinkulkarni/Documents/New project/Urban Waste System - India Ecosystem Lens.pptx"
MAP = "/Users/ashwinkulkarni/Downloads/codex-map-preview/Attachment287.pdf.png"


BG = RGBColor(247, 248, 244)
TEXT = RGBColor(34, 34, 34)
MUTED = RGBColor(95, 95, 95)
LIGHT = RGBColor(255, 255, 255)
BORDER = RGBColor(206, 210, 204)
RED = RGBColor(206, 66, 66)
ORANGE = RGBColor(230, 126, 34)
PURPLE = RGBColor(155, 89, 182)
BLUE = RGBColor(52, 152, 219)
GREEN = RGBColor(46, 204, 113)
GRAY = RGBColor(127, 140, 141)


def add_round_box(slide, x, y, w, h, fill=LIGHT, line=BORDER, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.2)
    return shp


def set_text_frame(tf, text="", font_size=20, bold=False, color=TEXT, align=PP_ALIGN.LEFT, margin=0.08):
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Aptos"
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    p.alignment = align
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    return p


def add_title(slide, title, subtitle=None):
    box = add_round_box(slide, Inches(0.2), Inches(0.15), Inches(12.9), Inches(0.75))
    set_text_frame(box.text_frame, title, font_size=22, bold=True)
    if subtitle:
        tx = slide.shapes.add_textbox(Inches(0.33), Inches(0.62), Inches(12.2), Inches(0.25))
        set_text_frame(tx.text_frame, subtitle, font_size=8, color=MUTED, margin=0)


def add_stat_card(slide, x, y, w, h, big, small, color):
    card = add_round_box(slide, x, y, w, h, fill=LIGHT, line=color)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = big
    r1.font.name = "Aptos"
    r1.font.size = Pt(18)
    r1.font.bold = True
    r1.font.color.rgb = color
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = small
    r2.font.name = "Aptos"
    r2.font.size = Pt(8)
    r2.font.color.rgb = MUTED


def add_bullet_card(slide, x, y, w, h, title, bullets, accent):
    card = add_round_box(slide, x, y, w, h, fill=LIGHT, line=accent)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = accent
    for bullet in bullets:
        pb = tf.add_paragraph()
        pb.level = 0
        pb.bullet = True
        run = pb.add_run()
        run.text = bullet
        run.font.name = "Aptos"
        run.font.size = Pt(9)
        run.font.color.rgb = TEXT


def add_label_card(slide, x, y, w, h, title, body, accent):
    card = add_round_box(slide, x, y, w, h, fill=LIGHT, line=accent)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Aptos"
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = accent
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Aptos"
    r2.font.size = Pt(8)
    r2.font.color.rgb = TEXT


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_title(
        slide,
        "Slide 1 | India-Wide Ecosystem Lens: What Exists, What’s Missing, Gaps & Bottlenecks",
        "Lens date: April 2, 2026 | Uses CPCB FY 2021-22 mass-balance plus GFC 2024-25 and PIB January/March 2026 operating signals",
    )
    add_stat_card(slide, Inches(0.3), Inches(1.0), Inches(3.0), Inches(0.75), "1,62,000 TPD", "Urban MSW currently produced in Indian cities", PURPLE)
    add_stat_card(slide, Inches(3.45), Inches(1.0), Inches(3.0), Inches(0.75), "60,478 vs 39,648", "Wards with >90% collection vs >90% source segregation", RED)
    add_stat_card(slide, Inches(6.6), Inches(1.0), Inches(3.0), Inches(0.75), "54%", "National processing rate in CPCB FY 2021-22 mass-balance", BLUE)
    add_stat_card(slide, Inches(9.75), Inches(1.0), Inches(3.0), Inches(0.75), "~25 crore MT", "Legacy waste identified; >62% processed by January 31, 2026", GRAY)

    add_bullet_card(
        slide, Inches(0.3), Inches(1.95), Inches(4.0), Inches(2.05),
        "What is already there",
        [
            "ULB-led collection backbone with concessionaires, fleets and ward systems",
            "A large informal recovery economy of waste pickers, kabadiwalas and scrap traders",
            "Growing MRF, compost, CBG, recycling and EPR infrastructure",
            "SBM-U 2.0, GFC ratings and dumpsite remediation momentum",
        ],
        GREEN,
    )
    add_bullet_card(
        slide, Inches(4.55), Inches(1.95), Inches(4.0), Inches(2.05),
        "What is not fully there yet",
        [
            "Consistent source segregation across households and bulk generators",
            "End-to-end traceability for low-value waste streams",
            "Reliable offtake for sorted material across geographies",
            "Formal inclusion of informal workers at system scale",
        ],
        ORANGE,
    )
    add_bullet_card(
        slide, Inches(8.8), Inches(1.95), Inches(4.2), Inches(2.05),
        "Structural gaps",
        [
            "Feedstock quality, contamination control and route discipline",
            "Working capital and project finance for aggregation and processing",
            "Data interoperability across ULBs, EPR managers and recyclers",
            "Scientific landfill only for true rejects, not routine overflow",
        ],
        PURPLE,
    )
    add_bullet_card(
        slide, Inches(0.3), Inches(4.2), Inches(8.2), Inches(1.95),
        "Binding bottlenecks on the current map",
        [
            "Segregation gap at source",
            "Post-collection leakage between collection and processing",
            "MRF quality / recovery bottlenecks",
            "Offtake bottlenecks between aggregators and recyclers",
            "Residual dependence on landfill / RDF / processing rejects",
            "Legacy dumpsites still absorbing land, money and management attention",
        ],
        RED,
    )
    img = slide.shapes.add_picture(MAP, Inches(8.7), Inches(4.15), width=Inches(4.15))
    img.line.color.rgb = BORDER
    note = add_round_box(slide, Inches(8.7), Inches(6.25), Inches(4.15), Inches(0.7), fill=LIGHT, line=BORDER)
    set_text_frame(
        note.text_frame,
        "This ecosystem view is additive, not exclusionary: public systems, entrepreneurs, for-profits, nonprofits, cooperatives, consultants, family-backed vehicles and informal actors all matter.",
        font_size=8,
        color=MUTED,
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_title(
        slide,
        "Slide 2 | Actors, Partners & Illustrative India-Wide Examples by System Block",
        "Illustrative, not exhaustive. Each block can host public, private, nonprofit, cooperative, advisory and informal actors.",
    )
    # left grid 3x3
    cards = [
        ("Consumers & Bulk Waste Generators", "Households, RWAs, campuses, hotels, malls, hospitals, stations, airports, factories. Examples: residential societies, Indian Railways sites, airports, IT parks, schools, hospitality chains.", PURPLE),
        ("Formal Waste Collectors", "ULBs, sanitation departments, ward contractors, concessionaires, fleets, worker teams. Examples: municipal corporations, Antony Waste, city-level collection operators.", BLUE),
        ("Informal Waste Collectors", "Waste pickers, itinerant buyers, kabadiwalas, cooperatives, worker groups. Examples: SWaCH, Hasiru Dala, Chintan / Safai Sena, collector networks linked to Plastics for Change.", ORANGE),
        ("MRFs / Sorting Centres", "Municipal MRFs, DWCCs, zero-waste operators, sorting lines, social enterprises. Examples: Saahas Zero Waste, Hasiru Dala Innovations, Chintan micro-MRFs, city MRF networks.", GREEN),
        ("Scrap Dealers", "Kabadi shops, balers, dry-waste godowns, mandi traders and small aggregators. Examples: local kabadi networks, neighborhood dry-waste traders, TheKabadiwala interfaces.", PURPLE),
        ("Aggregators", "Digital marketplaces, traceability platforms, EPR managers, material consolidators. Examples: Recykal, GEM Enviro, Plastics for Change, Lucro collection network, Cirqlo.", BLUE),
        ("Recyclers", "Plastic, metal, paper, e-waste, battery and multi-material recyclers. Examples: Banyan Nation, Lucro, Attero, Gravita, GEM Enviro, Deluxe.", ORANGE),
        ("Waste Transporters", "Municipal fleets, concessionaire fleets, micro-haulers, bale logistics and reverse logistics partners. Examples: ULB fleets, Antony fleets, producer take-back logistics, local transporters.", GREEN),
        ("Waste Processors", "Composters, biomethanation, CBG, RDF, WtE, co-processing and scientific landfill operators. Examples: GPS Renewables, IGRPL, municipal composters, large processing operators.", RED),
    ]
    x0, y0, w, h, gapx, gapy = Inches(0.28), Inches(1.0), Inches(2.7), Inches(1.38), Inches(0.15), Inches(0.15)
    idx = 0
    for r in range(3):
        for c in range(3):
            title, body, color = cards[idx]
            add_label_card(slide, x0 + c * (w + gapx), y0 + r * (h + gapy), w, h, title, body, color)
            idx += 1
    # right map and cross-cutting box
    map_box = add_round_box(slide, Inches(8.75), Inches(1.0), Inches(4.3), Inches(4.35), fill=LIGHT, line=BORDER)
    slide.shapes.add_picture(MAP, Inches(8.92), Inches(1.15), width=Inches(3.95))
    caption = slide.shapes.add_textbox(Inches(8.95), Inches(5.0), Inches(3.8), Inches(0.28))
    set_text_frame(caption.text_frame, "Same system map, now read as an India-wide actor landscape.", font_size=8, color=MUTED, margin=0)
    add_bullet_card(
        slide, Inches(8.75), Inches(5.45), Inches(4.3), Inches(1.65),
        "Cross-cutting partners that show up across blocks",
        [
            "MoHUA / SBM-U 2.0, SPCBs / PCCs, ULBs and state urban departments",
            "UNDP India, civil society groups, consultants, research institutions and accelerators",
            "Producer / brand EPR systems, marketplaces, family-backed and strategic capital, and entrepreneur-led platforms",
        ],
        GRAY,
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_title(
        slide,
        "Slide 3 | Capital Layer Attached to the Map: Funders, Philanthropy, CSR, Investors & Enablers",
        "Read capital as sitting on different blocks of the same system. This lens is intentionally inclusive: public, philanthropic, CSR, strategic, venture and family-backed capital all play a role.",
    )

    map_box = add_round_box(slide, Inches(0.3), Inches(1.0), Inches(4.15), Inches(5.65), fill=LIGHT, line=BORDER)
    slide.shapes.add_picture(MAP, Inches(0.48), Inches(1.18), width=Inches(3.8))
    caption = slide.shapes.add_textbox(Inches(0.48), Inches(5.15), Inches(3.75), Inches(1.2))
    set_text_frame(
        caption.text_frame,
        "Capital does not sit in one box. It shows up differently across first mile, inclusion, aggregation, recycling, processing and legacy remediation.",
        font_size=8,
        color=MUTED,
        margin=0,
    )

    capital_cards = [
        ("Consumers / source segregation", "ULB budgets, SBM-U IEC, CSR behavior-change campaigns and donor pilots. Examples: HUL-MoHUA, Tata Trusts, UNDP India.", PURPLE),
        ("Informal collectors", "Livelihood grants, worker inclusion philanthropy, fair-trade plastic models and CSR skilling. Examples: Body Shop x Plastics for Change, worker-support platforms.", ORANGE),
        ("Formal collection / transport", "Municipal opex, concessionaire finance, fleet leasing and producer take-back budgets for first-mile and reverse logistics.", BLUE),
        ("MRFs / sorting centres", "SBM-U capex, city PPPs, CSR pilots and blended finance for quality upgrades, digitization and traceability. Examples: Social Alpha, Startup Gateway pilots.", GREEN),
        ("Scrap dealers / aggregators", "Working capital, EPR manager spend, inventory finance and angel / family-backed capital for consolidation and route discipline.", PURPLE),
        ("Recyclers", "VC, growth equity, strategic investors, family offices and brand offtake. Examples: Circulate Capital, Claypond Capital, strategic corporate buyers.", BLUE),
        ("Processors", "Project finance, infrastructure funds, strategic energy investors and carbon-linked capital for compost, CBG, RDF and scientific processing.", RED),
        ("Landfill / legacy waste", "Public remediation funds, state support, climate / blended finance and technical assistance. Priority use: dumpsites and scientific rejects management.", GRAY),
    ]
    x0, y0, w, h = Inches(4.7), Inches(1.0), Inches(4.0), Inches(1.2)
    idx = 0
    for r in range(4):
        for c in range(2):
            title, body, color = capital_cards[idx]
            add_label_card(slide, x0 + c * Inches(4.2), y0 + r * Inches(1.3), w, h, title, body, color)
            idx += 1

    foot = add_round_box(slide, Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.3), fill=LIGHT, line=BORDER)
    set_text_frame(
        foot.text_frame,
        "Illustrative capital actors across the system include Tata Trusts, H&M Foundation, BMGF, USAID, UNDP India, ITC WOW, Coca-Cola India Foundation Anandana, Circulate Capital, Social Alpha, family offices and strategic investors.",
        font_size=8,
        color=MUTED,
    )

    prs.save(OUT)


if __name__ == "__main__":
    build()
